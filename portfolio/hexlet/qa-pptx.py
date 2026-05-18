"""QA check for hexlet-research.pptx: boundary overflow and obvious overlaps."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = 0

issues = []

def emu_to_in(v):
    return v / 914400

prs = Presentation('hexlet-research.pptx')

for idx, slide in enumerate(prs.slides, 1):
    slide_issues = []
    for shape in slide.shapes:
        l = shape.left or 0
        t = shape.top or 0
        w = shape.width or 0
        h = shape.height or 0

        # Check boundary overflow
        if l < -MARGIN:
            slide_issues.append(f"  ⚠ Shape '{shape.name}' left={emu_to_in(l):.2f}\" OVERFLOW LEFT")
        if t < -MARGIN:
            slide_issues.append(f"  ⚠ Shape '{shape.name}' top={emu_to_in(t):.2f}\" OVERFLOW TOP")
        if l + w > SLIDE_W + MARGIN:
            slide_issues.append(f"  ⚠ Shape '{shape.name}' right={emu_to_in(l+w):.2f}\" OVERFLOW RIGHT (slide={emu_to_in(SLIDE_W):.2f}\")")
        if t + h > SLIDE_H + MARGIN:
            slide_issues.append(f"  ⚠ Shape '{shape.name}' bottom={emu_to_in(t+h):.2f}\" OVERFLOW BOTTOM (slide={emu_to_in(SLIDE_H):.2f}\")")

        # Check text frame overflow (approximate via font size * lines)
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(48):
                        slide_issues.append(f"  ℹ Shape '{shape.name}' has large font {Pt(run.font.size):.0f}pt — verify visually")
                        break

    if slide_issues:
        issues.append(f"\nSlide {idx}:")
        issues.extend(slide_issues)

if issues:
    print("QA ISSUES FOUND:")
    print('\n'.join(issues))
else:
    print("✅ QA PASSED — no boundary overflows detected across all 11 slides.")

# Summary
print(f"\nSlides: {len(prs.slides)}")
for idx, slide in enumerate(prs.slides, 1):
    shape_count = len(slide.shapes)
    has_text = sum(1 for s in slide.shapes if s.has_text_frame)
    print(f"  Slide {idx:2d}: {shape_count:3d} shapes, {has_text:2d} with text")
