"""
Build egeland_cjm.pptx — CJM Parent Journey through EGE Land website.
Brand colors, sticky-note comments, screenshot layout.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

BASE = Path(__file__).parent
SCREENSHOTS = BASE / "screenshots"
OUT_FILE = BASE / "egeland_cjm.pptx"

# ── Brand colours ────────────────────────────────────────────────────────────
BG_HEX     = RGBColor(0x1A, 0x1A, 0x1A)
PURPLE_HEX = RGBColor(0x7C, 0x3A, 0xED)
LIME_HEX   = RGBColor(0xC8, 0xFF, 0x00)
WHITE_HEX  = RGBColor(0xFF, 0xFF, 0xFF)
STICKER_BG = RGBColor(0xFF, 0xF1, 0x76)
STICKER_FG = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_HEX   = RGBColor(0x88, 0x88, 0x88)

# ── Slide dimensions 16:9 ────────────────────────────────────────────────────
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG_HEX):
    bg_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()


def txb(slide, text, x, y, w, h, size=18, bold=False, color=WHITE_HEX,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Arial"
    return box


def step_number(slide, num):
    """Lime step badge top-left."""
    txb(slide, f"Шаг {num}", Cm(1), Cm(0.4), Cm(6), Cm(1.2),
        size=20, bold=True, color=LIME_HEX)


def add_image(slide, img_path, x, y, w, h):
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), x, y, w, h)


def sticker(slide, lines, x, y, w):
    """Yellow sticky-note card with multiple lines."""
    line_h = Cm(1.1)
    padding_v = Cm(0.6)
    n = len(lines)
    h = n * line_h + 2 * padding_v

    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = STICKER_BG
    shape.line.color.rgb = RGBColor(0xDD, 0xCC, 0x00)
    shape.line.width = Pt(1)

    # Adjust corner radius via XML directly
    try:
        from pptx.oxml.ns import qn
        spPr = shape.element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", "roundRect")
    except Exception:
        pass

    tf = shape.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4) if i > 0 else Pt(8)
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(14)
        run.font.bold  = False
        run.font.color.rgb = STICKER_FG
        run.font.name  = "Arial"

    return shape


def accent_bar(slide):
    """Thin lime bar at top."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Cm(0.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIME_HEX
    bar.line.fill.background()


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s)

# Decorative purple rectangle left
rect = s.shapes.add_shape(1, 0, 0, Cm(0.8), SLIDE_H)
rect.fill.solid()
rect.fill.fore_color.rgb = PURPLE_HEX
rect.line.fill.background()

# Lime horizontal accent
accent_bar(s)

txb(s, "CJM: путь родителя", Cm(2.5), Cm(4), Cm(29), Cm(3.5),
    size=44, bold=True, color=WHITE_HEX)

txb(s, "ЕГЭ Лэнд", Cm(2.5), Cm(7.3), Cm(15), Cm(1.8),
    size=36, bold=True, color=LIME_HEX)

txb(s, "Исследование воронки перед откликом на вакансию PMM/Growth",
    Cm(2.5), Cm(9.4), Cm(28), Cm(2),
    size=20, color=GRAY_HEX, italic=True)

txb(s, "Надя Романенко", Cm(2.5), Cm(16.5), Cm(20), Cm(1.5),
    size=16, color=GRAY_HEX)

# Decorative lime circle top-right
circ = s.shapes.add_shape(9, Cm(27), Cm(1), Cm(6), Cm(6))  # oval
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x20, 0x10, 0x40)
circ.line.color.rgb = PURPLE_HEX
circ.line.width = Pt(2)

txb(s, "egeland.ru", Cm(26.5), Cm(2.5), Cm(7), Cm(2),
    size=16, color=PURPLE_HEX, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# Helper: content slide with screenshot left + stickers right
# ────────────────────────────────────────────────────────────────────────────
IMG_X = Cm(0.6)
IMG_Y = Cm(1.8)
IMG_W = Cm(19)
IMG_H = Cm(16.5)

STICKER_X = Cm(20.5)
STICKER_W = Cm(12.5)


def content_slide(step_num, title, img_file, comments):
    s = add_slide()
    bg(s)
    accent_bar(s)
    step_number(s, step_num)

    txb(s, title, Cm(7), Cm(0.3), Cm(25), Cm(1.4),
        size=18, bold=True, color=WHITE_HEX)

    add_image(s, SCREENSHOTS / img_file, IMG_X, IMG_Y, IMG_W, IMG_H)

    y = IMG_Y
    for comment_lines in comments:
        sticker(s, comment_lines, STICKER_X, y, STICKER_W)
        # Estimate height
        h = len(comment_lines) * Cm(1.1) + Cm(1.2)
        y += h + Cm(0.5)

    return s


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Landing
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=1,
    title="Лендинг — вкладка «Родитель»",
    img_file="01_landing.png",
    comments=[
        ["✅ Оффер есть, УТП есть"],
        ["⚠️ «Выбери класс» — смешение ролей:",
         "я во вкладке Родитель, но форма",
         "спрашивает про меня, а не про ребёнка"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Role selection
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=2,
    title="Шаг «Кто вы» — выбор роли",
    img_file="02_role_selection.png",
    comments=[
        ["Воронка разделяет ученика и родителя —",
         "но ведёт их по одинаковому пути.",
         "В чём тогда смысл разделения?"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Class selection
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=3,
    title="Шаг «Класс» — повторный вопрос",
    img_file="03_class_selection.png",
    comments=[
        ["⚠️ Повторяется вопрос — класс уже",
         "спрашивали на лендинге"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Subjects
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=4,
    title="Шаг «Предметы»",
    img_file="04_subjects.png",
    comments=[
        ["Кнопка «Узнать подробности»",
         "перебрасывает на бота — кажется,",
         "там отвечают реальные люди.",
         "Это не очевидно пользователю"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Upsell
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=5,
    title="Апселл — третий предмет",
    img_file="05_upsell.png",
    comments=[
        ["💡 Механика увеличения среднего чека —",
         "третий предмет бесплатно летом. Грамотно"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Tariff Standard
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=6,
    title="Тариф Стандарт",
    img_file="06_tariff_standard.png",
    comments=[
        ["❓ Незакрытые вопросы:"],
        ["Кем составляется индивидуальная",
         "траектория?"],
        ["Как выглядит помощь с домашками?"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Tariff Premium
# ────────────────────────────────────────────────────────────────────────────
content_slide(
    step_num=7,
    title="Тариф Премиум",
    img_file="07_tariff_premium.png",
    comments=[
        ["❓ Ускоренная проверка — а сколько ждать",
         "в Стандарте?"],
        ["Мини-группа до 20 — сколько в обычной?"],
        ["Как выглядит поддержка психолога?"],
    ]
)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Summary (no screenshot)
# ────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)

txb(s, "Что я вижу как PMM", Cm(2), Cm(0.5), Cm(30), Cm(2),
    size=34, bold=True, color=WHITE_HEX)

insights = [
    ("Воронка технически разделена на две ЦА — но не работает с их разными страхами",
     Cm(2), Cm(3.8)),
    ("6+ незакрытых вопросов у родителя до оплаты — это и есть потеря денег",
     Cm(2), Cm(8.2)),
    ("Знаю как это чинить: CJM → гипотезы → цифры. Кейсы прилагаются",
     Cm(2), Cm(12.5)),
]

for i, (text, x, y) in enumerate(insights):
    # Lime accent dot / number
    dot = s.shapes.add_shape(9, x, y + Cm(0.35), Cm(0.7), Cm(0.7))
    dot.fill.solid()
    dot.fill.fore_color.rgb = LIME_HEX
    dot.line.fill.background()

    # Block background
    block = s.shapes.add_shape(1, x + Cm(1.1), y, Cm(28.5), Cm(3.2))
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor(0x28, 0x18, 0x48)
    block.line.color.rgb = PURPLE_HEX
    block.line.width = Pt(1.5)

    # Text inside block
    tf = block.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_before = Pt(10)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = False
    run.font.color.rgb = WHITE_HEX
    run.font.name = "Arial"

    # Lime number label
    txb(s, str(i + 1), x + Cm(0.13), y + Cm(0.2), Cm(0.5), Cm(0.7),
        size=14, bold=True, color=BG_HEX, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"Saved: {OUT_FILE}")
